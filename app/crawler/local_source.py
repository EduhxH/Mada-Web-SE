import hashlib
import io
import json
import logging
import zipfile
from collections import Counter
from dataclasses import dataclass, field
from pathlib import Path

import docx
from bs4 import BeautifulSoup
from pptx import Presentation
from pypdf import PdfReader

from urllib.parse import unquote, urlparse

from app.indexing.tokenizer import remover_acentos
from app.models.document import Documento

logging.getLogger("pypdf").setLevel(logging.ERROR)

EXTENSOES_TEXTO = {".txt", ".md", ".cs"}
EXTENSOES_HTML = {".html", ".htm"}
EXTENSOES_SUPORTADAS = (
    EXTENSOES_TEXTO | EXTENSOES_HTML | {".pdf", ".docx", ".pptx"}
)

ETIQUETAS_SEM_CONTEUDO = ("script", "style", "nav", "header", "footer", "noscript")

PASTAS_IGNORADAS = ("bin/", "obj/", ".vs/", "packages/", "__pycache__/")

PADROES_PRIVADOS = ("notas", "pauta", "classifica")

# Ficheiros do proprio sistema: segredos, codigos de acesso, bases de dados.
# Nunca devem entrar no indice, aconteca o que acontecer ao caminho indicado.
NOMES_PROIBIDOS = frozenset(
    {"segredo.txt", "participantes.json", ".env", ".env.example"}
)
EXTENSOES_PROIBIDAS = {".sqlite3", ".sqlite", ".db", ".key", ".pem"}

NOME_MANIFESTO = "_origens.json"

MOTIVO_PRIVADO = "possivel dado pessoal"
MOTIVO_SISTEMA = "ficheiro do sistema"
MOTIVO_FORMATO = "formato nao suportado"
MOTIVO_BUILD = "artefacto de build"
MOTIVO_SEM_TEXTO = "sem texto extraivel"
MOTIVO_CORROMPIDO = "ficheiro ilegivel"
MOTIVO_ZIP_ANINHADO = "zip dentro de zip"
MOTIVO_DUPLICADO = "documento repetido"

# 48 bits: com milhares de documentos a hipotese de colisao e de 1 em
# centenas de milhoes, e a origem ja e unica por documento.
BYTES_ID = 6


@dataclass
class Relatorio:
    por_disciplina: Counter = field(default_factory=Counter)
    por_formato: Counter = field(default_factory=Counter)
    ignorados: list[tuple[str, str]] = field(default_factory=list)

    def ignorar(self, nome: str, motivo: str) -> None:
        self.ignorados.append((nome, motivo))

    def motivos(self) -> Counter:
        return Counter(motivo for _, motivo in self.ignorados)

    def por_motivo(self, motivo: str) -> list[str]:
        return [nome for nome, m in self.ignorados if m == motivo]


def carregar(caminho: str | Path) -> tuple[list[Documento], Relatorio]:
    raiz = Path(caminho)
    relatorio = Relatorio()
    documentos: list[Documento] = []

    if raiz.is_dir():
        arquivos = sorted(p for p in raiz.rglob("*") if p.is_file())
    else:
        arquivos = [raiz]

    manifestos: dict[Path, dict[str, str]] = {}
    vistos: set[int] = set()
    for arquivo in arquivos:
        if arquivo.name == NOME_MANIFESTO:
            continue
        pasta = arquivo.parent
        if pasta not in manifestos:
            manifestos[pasta] = _ler_manifesto(pasta)
        entrada = manifestos[pasta].get(arquivo.name)
        if isinstance(entrada, dict):
            origem_do_ficheiro = entrada.get("url") or str(arquivo)
            titulo_do_manifesto = entrada.get("titulo") or ""
        else:
            origem_do_ficheiro = entrada or str(arquivo)
            titulo_do_manifesto = ""
        _processar(
            nome=arquivo.name,
            dados=arquivo.read_bytes(),
            origem=origem_do_ficheiro,
            titulo_dado=titulo_do_manifesto,
            disciplina=_disciplina(raiz, arquivo),
            documentos=documentos,
            relatorio=relatorio,
            vistos=vistos,
        )

    return documentos, relatorio


def id_estavel(origem: str) -> int:
    """Id derivado da origem, nao da ordem de leitura.

    Com ids sequenciais, acrescentar um ficheiro no inicio da pasta
    deslocava todos os outros - e os cliques ja registados passavam a
    apontar para documentos errados. Assim, o mesmo ficheiro tem sempre o
    mesmo id, reindexe-se as vezes que for preciso.
    """
    resumo = hashlib.blake2b(
        origem.encode("utf-8"), digest_size=BYTES_ID
    ).hexdigest()
    return int(resumo, 16)


def _ler_manifesto(pasta: Path) -> dict[str, str]:
    caminho = pasta / NOME_MANIFESTO
    if not caminho.exists():
        return {}
    try:
        dados = json.loads(caminho.read_text(encoding="utf-8"))
    except (json.JSONDecodeError, OSError):
        return {}
    return dados if isinstance(dados, dict) else {}


def titulo_de_url(url: str) -> str:
    """Titulo legivel a partir do URL, em vez do nome do ficheiro guardado.

    .../uploads/REGULAMENTO-INTERNO-APROVADO.pdf -> "REGULAMENTO INTERNO APROVADO"
    """
    caminho = unquote(urlparse(url).path).rstrip("/")
    if not caminho:
        return ""
    nome = Path(caminho).name
    for extensao in (".pdf", ".docx", ".pptx", ".html", ".htm"):
        if nome.lower().endswith(extensao):
            nome = nome[: -len(extensao)]
            break
    nome = nome.replace("-", " ").replace("_", " ")
    return " ".join(nome.split())


def _disciplina(raiz: Path, arquivo: Path) -> str:
    if not raiz.is_dir():
        return raiz.parent.name
    relativo = arquivo.relative_to(raiz)
    return relativo.parts[-2] if len(relativo.parts) > 1 else raiz.name


def _e_privado(nome: str) -> bool:
    normalizado = remover_acentos(nome.lower())
    return any(padrao in normalizado for padrao in PADROES_PRIVADOS)


def _e_artefacto(caminho_interno: str) -> bool:
    normalizado = caminho_interno.replace("\\", "/").lower()
    return any(pasta in normalizado for pasta in PASTAS_IGNORADAS)


def _processar(
    nome: str,
    dados: bytes,
    origem: str,
    disciplina: str,
    documentos: list[Documento],
    relatorio: Relatorio,
    titulo_dado: str = "",
    vistos: set[int] | None = None,
    dentro_de_zip: bool = False,
) -> None:
    vistos = vistos if vistos is not None else set()
    extensao = Path(nome).suffix.lower()

    if nome.lower() in NOMES_PROIBIDOS or extensao in EXTENSOES_PROIBIDAS:
        relatorio.ignorar(origem, MOTIVO_SISTEMA)
        return

    if _e_privado(nome):
        relatorio.ignorar(origem, MOTIVO_PRIVADO)
        return

    if extensao == ".zip":
        if dentro_de_zip:
            relatorio.ignorar(origem, MOTIVO_ZIP_ANINHADO)
            return
        _processar_zip(dados, origem, disciplina, documentos, relatorio, vistos)
        return

    if extensao not in EXTENSOES_SUPORTADAS:
        relatorio.ignorar(origem, MOTIVO_FORMATO)
        return

    try:
        unidades = _extrair(extensao, dados)
    except Exception:
        relatorio.ignorar(origem, MOTIVO_CORROMPIDO)
        return

    if not unidades:
        relatorio.ignorar(origem, MOTIVO_SEM_TEXTO)
        return

    # Um titulo do manifesto (nome do recurso no Moodle) vale mais que o
    # nome do ficheiro guardado ou o ultimo segmento do URL.
    base = titulo_dado or Path(nome).stem
    if not titulo_dado and origem.startswith(("http://", "https://")):
        base = titulo_de_url(origem) or base
    if extensao in EXTENSOES_HTML:
        origem = origem_declarada(dados) or origem
        base = titulo_html(dados) or base

    guardados = 0
    for conteudo, sufixo_titulo, sufixo_origem in unidades:
        origem_completa = origem + sufixo_origem
        identificador = id_estavel(origem_completa)
        if identificador in vistos:
            relatorio.ignorar(origem_completa, MOTIVO_DUPLICADO)
            continue
        vistos.add(identificador)
        documentos.append(
            Documento(
                id=identificador,
                titulo=base + sufixo_titulo,
                texto=conteudo,
                origem=origem_completa,
                disciplina=disciplina,
            )
        )
        guardados += 1
    if not guardados:
        return
    relatorio.por_disciplina[disciplina] += guardados
    relatorio.por_formato[extensao] += 1


def _processar_zip(
    dados: bytes,
    origem: str,
    disciplina: str,
    documentos: list[Documento],
    relatorio: Relatorio,
    vistos: set[int],
) -> None:
    try:
        arquivo_zip = zipfile.ZipFile(io.BytesIO(dados))
    except zipfile.BadZipFile:
        relatorio.ignorar(origem, MOTIVO_CORROMPIDO)
        return

    with arquivo_zip:
        for entrada in arquivo_zip.infolist():
            if entrada.is_dir():
                continue
            interno = entrada.filename
            rotulo = f"{origem}!{interno}"
            if _e_artefacto(interno):
                relatorio.ignorar(rotulo, MOTIVO_BUILD)
                continue
            try:
                conteudo = arquivo_zip.read(entrada)
            except Exception:
                relatorio.ignorar(rotulo, MOTIVO_CORROMPIDO)
                continue
            _processar(
                nome=Path(interno).name,
                dados=conteudo,
                origem=rotulo,
                disciplina=disciplina,
                documentos=documentos,
                relatorio=relatorio,
                vistos=vistos,
                dentro_de_zip=True,
            )


def _extrair(extensao: str, dados: bytes) -> list[tuple[str, str, str]]:
    if extensao in EXTENSOES_HTML:
        return _unidades_html(dados)
    if extensao == ".pdf":
        return _unidades_pdf(dados)
    if extensao == ".pptx":
        return _unidades_pptx(dados)
    if extensao == ".docx":
        return _unidades_docx(dados)
    return _unidades_texto(dados)


def _unidades_pdf(dados: bytes) -> list[tuple[str, str, str]]:
    leitor = PdfReader(io.BytesIO(dados))
    unidades = []
    for numero, pagina in enumerate(leitor.pages, start=1):
        texto = pagina.extract_text() or ""
        if texto.strip():
            unidades.append((texto, f" - pagina {numero}", f"#pagina={numero}"))
    return unidades


def _unidades_pptx(dados: bytes) -> list[tuple[str, str, str]]:
    apresentacao = Presentation(io.BytesIO(dados))
    unidades = []
    for numero, slide in enumerate(apresentacao.slides, start=1):
        partes = [
            forma.text_frame.text
            for forma in slide.shapes
            if forma.has_text_frame
        ]
        texto = "\n".join(partes)
        if texto.strip():
            unidades.append((texto, f" - slide {numero}", f"#slide={numero}"))
    return unidades


def _unidades_docx(dados: bytes) -> list[tuple[str, str, str]]:
    documento = docx.Document(io.BytesIO(dados))
    partes = [paragrafo.text for paragrafo in documento.paragraphs]
    for tabela in documento.tables:
        for linha in tabela.rows:
            partes.append(" ".join(celula.text for celula in linha.cells))
    texto = "\n".join(partes)
    return [(texto, "", "")] if texto.strip() else []


def _unidades_texto(dados: bytes) -> list[tuple[str, str, str]]:
    texto = dados.decode("utf-8", errors="replace")
    return [(texto, "", "")] if texto.strip() else []


def _unidades_html(dados: bytes) -> list[tuple[str, str, str]]:
    sopa = BeautifulSoup(dados.decode("utf-8", errors="replace"), "html.parser")
    for etiqueta in sopa(list(ETIQUETAS_SEM_CONTEUDO)):
        etiqueta.decompose()
    texto = sopa.get_text(" ", strip=True)
    return [(texto, "", "")] if texto.strip() else []


def origem_declarada(dados: bytes) -> str | None:
    sopa = BeautifulSoup(dados.decode("utf-8", errors="replace"), "html.parser")
    marca = sopa.find("meta", attrs={"name": "madalena-origem"})
    return marca.get("content") if marca else None


def titulo_html(dados: bytes) -> str | None:
    sopa = BeautifulSoup(dados.decode("utf-8", errors="replace"), "html.parser")
    if sopa.title and sopa.title.string:
        return sopa.title.string.strip()
    cabecalho = sopa.find("h1")
    return cabecalho.get_text(strip=True) if cabecalho else None
